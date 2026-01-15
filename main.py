from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import uvicorn
import os
import uuid
from pathlib import Path
import PyPDF2
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import json
from datetime import datetime
import requests as http_requests
from pptx import Presentation
from PIL import Image
import pytesseract
import io

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
app = FastAPI(title="RAG Chatbot API")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

# Initialize embedding model
print("Loading embedding model...")
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
EMBEDDING_DIM = 384
print("Embedding model loaded!")

# HuggingFace API configuration (FREE!)
HF_API_URL = "https://api-inference.huggingface.co/models/mistralai/Mistral-7B-Instruct-v0.2"
HF_API_TOKEN = os.getenv("HF_TOKEN", "")

BASE_DIR = Path("rag_data")
SUBJECTS_DIR = BASE_DIR / "subjects"
BASE_DIR.mkdir(exist_ok=True)
SUBJECTS_DIR.mkdir(exist_ok=True)

# In-memory storage (use database in production)
subjects_db = {}
documents_db = {}

def get_subject_folder(subject_id: str) -> Path:
    """Get or create a dedicated folder for a subject"""
    subject_folder = SUBJECTS_DIR / subject_id
    subject_folder.mkdir(exist_ok=True)
    
    # Create subfolders
    (subject_folder / "documents").mkdir(exist_ok=True)
    (subject_folder / "vector_db").mkdir(exist_ok=True)
    (subject_folder / "metadata").mkdir(exist_ok=True)
    
    return subject_folder

def get_subject_info_path(subject_id: str) -> Path:
    """Get path to subject info JSON file"""
    return get_subject_folder(subject_id) / "metadata" / "subject_info.json"

def save_subject_info(subject_id: str, info: dict):
    """Save subject information to disk"""
    info_path = get_subject_info_path(subject_id)
    with open(info_path, 'w') as f:
        json.dump(info, f, indent=2)

def load_subject_info(subject_id: str) -> dict:
    """Load subject information from disk"""
    info_path = get_subject_info_path(subject_id)
    if info_path.exists():
        with open(info_path, 'r') as f:
            return json.load(f)
    return None

# Models
class Subject(BaseModel):
    name: str

class ChatRequest(BaseModel):
    subject_id: str
    question: str
    top_k: int = 3

class ChatResponse(BaseModel):
    answer: str
    sources: List[str]

# Helper functions
def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file"""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file including images using OCR"""
    text = ""
    
    try:
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n=== Slide {slide_num} ===\n"
            
            # Extract text from text boxes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        text += row_text + "\n"
                
                # Extract text from images using OCR
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Convert to PIL Image
                        img = Image.open(io.BytesIO(image_bytes))
                        
                        # Perform OCR
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            text += f"\n[Image text]: {ocr_text}\n"
                    except Exception as e:
                        print(f"  ⚠ Could not extract text from image: {e}")
                        continue
            
            text += "\n"
        
        print(f"  ✓ Extracted text from {len(prs.slides)} slides")
        return text
        
    except Exception as e:
        print(f"  ✗ Error processing PowerPoint: {e}")
        return ""

def extract_text_from_image(file_path: str) -> str:
    """Extract text from image files using OCR"""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        print(f"  ✓ Extracted text from image using OCR")
        return text
    except Exception as e:
        print(f"  ✗ Error processing image: {e}")
        return ""

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """Split text into overlapping chunks"""
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += (chunk_size - overlap)
    
    return chunks

def load_vector_db(subject_id: str):
    """Load FAISS index and metadata for a subject"""
    subject_folder = get_subject_folder(subject_id)
    index_path = subject_folder / "vector_db" / "faiss.index"
    meta_path = subject_folder / "vector_db" / "metadata.json"
    
    if not index_path.exists():
        return None, None
    
    index = faiss.read_index(str(index_path))
    with open(meta_path, 'r') as f:
        metadata = json.load(f)
    
    return index, metadata

def save_vector_db(subject_id: str, index, metadata):
    """Save FAISS index and metadata for a subject"""
    subject_folder = get_subject_folder(subject_id)
    index_path = subject_folder / "vector_db" / "faiss.index"
    meta_path = subject_folder / "vector_db" / "metadata.json"
    
    faiss.write_index(index, str(index_path))
    with open(meta_path, 'w') as f:
        json.dump(metadata, f, indent=2)

def generate_answer(question: str, context: str) -> str:
    prompt = f"""<s>[INST] You are a helpful assistant. Answer the question based ONLY on the provided context. If the answer cannot be found in the context, say "No information found in the subject documents."

Context:
{context[:1000]}

Question: {question}

Provide a concise answer based only on the context above. [/INST]"""

    headers = {"Content-Type": "application/json"}
    if HF_API_TOKEN:
        headers["Authorization"] = f"Bearer {HF_API_TOKEN}"
    
    try:
        print("Calling HuggingFace API...")
        response = http_requests.post(
            HF_API_URL,
            headers=headers,
            json={
                "inputs": prompt,
                "parameters": {
                    "max_new_tokens": 200,
                    "temperature": 0.1,
                    "top_p": 0.95,
                    "return_full_text": False
                },
                "options": {
                    "wait_for_model": True
                }
            },
            timeout=60
        )
        
        print(f"API Response Status: {response.status_code}")
        
        if response.status_code == 200:
            result = response.json()
            print(f"API Response: {result}")
            
            if isinstance(result, list) and len(result) > 0:
                answer = result[0].get("generated_text", "").strip()
                if answer:
                    print(f"✓ Generated answer: {answer[:100]}...")
                    return answer
            elif isinstance(result, dict) and "generated_text" in result:
                answer = result["generated_text"].strip()
                if answer:
                    print(f"✓ Generated answer: {answer[:100]}...")
                    return answer
            
            print(f"Unexpected API response format: {result}")
        
        elif response.status_code == 503:
            print("Model is loading on HuggingFace servers, please wait...")
            return "The AI model is currently loading. Please try again in 10-20 seconds."
        
        else:
            print(f"API Error {response.status_code}: {response.text}")
            
    except http_requests.exceptions.Timeout:
        print("Request timed out")
        return "The request timed out. Please try again."
    except Exception as e:
        print(f"Error calling HuggingFace API: {type(e).__name__}: {e}")
    
    # If we get here, use simple extractive answer
    print("Falling back to extractive answer...")
    return extract_answer_from_context(question, context)

def extract_answer_from_context(question: str, context: str) -> str:
    """
    Simple extractive answer - finds the most relevant sentence
    """
    # Split context into sentences
    sentences = [s.strip() for s in context.replace('\n', ' ').split('.') if s.strip()]
    
    # Find most relevant sentence (simple keyword matching)
    question_words = set(question.lower().split())
    best_sentence = ""
    best_score = 0
    
    for sentence in sentences[:10]:  # Check first 10 sentences
        sentence_words = set(sentence.lower().split())
        score = len(question_words & sentence_words)
        if score > best_score:
            best_score = score
            best_sentence = sentence
    
    if best_sentence:
        return best_sentence + "."
    else:
        return sentences[0] + "." if sentences else "No specific answer found in the documents."

# API Endpoints
@app.post("/subjects")
async def create_subject(subject: Subject):
    """Create a new subject"""
    subject_id = str(uuid.uuid4())
    subject_info = {
        "id": subject_id,
        "name": subject.name,
        "created_at": datetime.now().isoformat(),
        "document_count": 0,
        "documents": []
    }
    
    subjects_db[subject_id] = subject_info
    
    # Save to disk
    save_subject_info(subject_id, subject_info)
    
    print(f"✓ Created subject: {subject.name} (ID: {subject_id[:8]}...)")
    print(f"  Folder: rag_data/subjects/{subject_id}/")
    
    return subject_info

@app.get("/subjects")
async def list_subjects():
    """List all subjects"""
    # Load subjects from disk if not in memory
    if not subjects_db:
        for subject_folder in SUBJECTS_DIR.iterdir():
            if subject_folder.is_dir():
                subject_id = subject_folder.name
                info = load_subject_info(subject_id)
                if info:
                    subjects_db[subject_id] = info
    
    return list(subjects_db.values())

@app.post("/subjects/{subject_id}/documents")
async def upload_documents(
    subject_id: str,
    files: List[UploadFile] = File(...)
):
    """Upload documents to a subject"""
    if subject_id not in subjects_db:
        # Try loading from disk
        info = load_subject_info(subject_id)
        if info:
            subjects_db[subject_id] = info
        else:
            raise HTTPException(status_code=404, detail="Subject not found")
    
    subject_folder = get_subject_folder(subject_id)
    documents_folder = subject_folder / "documents"
    
    # Load or create vector database
    index, metadata = load_vector_db(subject_id)
    if index is None:
        index = faiss.IndexFlatIP(EMBEDDING_DIM)  # Inner product (cosine similarity)
        metadata = {"chunks": [], "sources": [], "document_info": []}
    
    uploaded_files = []
    
    for file in files:
        # Create unique filename with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_filename = f"{timestamp}_{file.filename}"
        file_path = documents_folder / safe_filename
        
        # Save uploaded file
        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)
        
        print(f"  ✓ Saved: {safe_filename}")
        
        # Extract text based on file type
        text = ""
        if file.filename.lower().endswith('.pdf'):
            text = extract_text_from_pdf(str(file_path))
        elif file.filename.lower().endswith('.txt'):
            text = extract_text_from_txt(str(file_path))
        elif file.filename.lower().endswith(('.ppt', '.pptx')):
            text = extract_text_from_pptx(str(file_path))
        elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
            text = extract_text_from_image(str(file_path))
        else:
            print(f"  ✗ Skipped: {file.filename} (unsupported format)")
            continue
        
        if not text.strip():
            print(f"  ⚠ No text extracted from {file.filename}")
            continue
        
        # Chunk text
        chunks = chunk_text(text)
        print(f"  ✓ Extracted {len(chunks)} chunks from {file.filename}")
        
        # Generate embeddings
        embeddings = embedding_model.encode(chunks, normalize_embeddings=True)
        
        # Add to FAISS index
        index.add(embeddings.astype('float32'))
        
        # Store metadata
        doc_info = {
            "original_name": file.filename,
            "saved_name": safe_filename,
            "upload_time": datetime.now().isoformat(),
            "chunk_count": len(chunks),
            "file_size": len(content)
        }
        
        for chunk in chunks:
            metadata["chunks"].append(chunk)
            metadata["sources"].append(file.filename)
        
        metadata["document_info"].append(doc_info)
        
        uploaded_files.append(file.filename)
        
        # Update document count
        subjects_db[subject_id]["document_count"] += 1
        if "documents" not in subjects_db[subject_id]:
            subjects_db[subject_id]["documents"] = []
        subjects_db[subject_id]["documents"].append(doc_info)
    
    # Save vector database
    save_vector_db(subject_id, index, metadata)
    
    # Save updated subject info
    save_subject_info(subject_id, subjects_db[subject_id])
    
    print(f"✓ Upload complete for subject: {subjects_db[subject_id]['name']}")
    
    return {
        "subject_id": subject_id,
        "subject_name": subjects_db[subject_id]["name"],
        "uploaded_files": uploaded_files,
        "total_chunks": len(metadata["chunks"]),
        "document_count": subjects_db[subject_id]["document_count"]
    }

@app.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest):
    """Ask a question to a subject"""
    print(f"\n=== CHAT REQUEST ===")
    print(f"Subject ID: {request.subject_id}")
    print(f"Question: {request.question}")
    
    if request.subject_id not in subjects_db:
        raise HTTPException(status_code=404, detail="Subject not found")
    
    # Load vector database
    index, metadata = load_vector_db(request.subject_id)
    
    print(f"Index loaded: {index is not None}")
    if index:
        print(f"Total vectors in index: {index.ntotal}")
        print(f"Total chunks in metadata: {len(metadata.get('chunks', []))}")
    
    if index is None or index.ntotal == 0:
        print("ERROR: No vectors in database!")
        return ChatResponse(
            answer="No information found in the subject documents.",
            sources=[]
        )
    
    # Embed question
    question_embedding = embedding_model.encode(
        [request.question], 
        normalize_embeddings=True
    ).astype('float32')
    
    # Search similar chunks
    k = min(request.top_k, index.ntotal)
    distances, indices = index.search(question_embedding, k)
    
    print(f"Top {k} similarity scores: {distances[0]}")
    print(f"Top {k} chunk indices: {indices[0]}")
    
    # Relevance threshold (cosine similarity > 0.3)
    relevance_threshold = 0.2  # Lowered from 0.3
    
    # Collect relevant chunks
    relevant_chunks = []
    relevant_sources = []
    
    for dist, idx in zip(distances[0], indices[0]):
        print(f"Checking chunk {idx}: similarity={dist:.4f}, threshold={relevance_threshold}")
        if dist > relevance_threshold:
            chunk_text = metadata["chunks"][idx]
            print(f"✓ RELEVANT: {chunk_text[:100]}...")
            relevant_chunks.append(chunk_text)
            relevant_sources.append(metadata["sources"][idx])
        else:
            print(f"✗ Not relevant enough")
    
    if not relevant_chunks:
        print("ERROR: No chunks passed relevance threshold!")
        return ChatResponse(
            answer="No information found in the subject documents.",
            sources=[]
        )
    
    # Generate answer using context
    context = "\n\n".join(relevant_chunks)
    
    # Limit context size
    max_context_length = 1000
    if len(context) > max_context_length:
        context = context[:max_context_length] + "..."
    
    print(f"Context to send to LLM:\n{context[:200]}...\n")
    
    # Use HuggingFace API for answer generation
    answer = generate_answer(request.question, context)
    
    print(f"Final answer: {answer[:200]}...")
    
    return ChatResponse(
        answer=answer,
        sources=list(set(relevant_sources))
    )

@app.get("/subjects/{subject_id}/details")
async def get_subject_details(subject_id: str):
    """Get detailed information about a subject including all documents"""
    if subject_id not in subjects_db:
        info = load_subject_info(subject_id)
        if info:
            subjects_db[subject_id] = info
        else:
            raise HTTPException(status_code=404, detail="Subject not found")
    
    subject_folder = get_subject_folder(subject_id)
    documents_folder = subject_folder / "documents"
    
    # List all documents in folder
    documents = []
    if documents_folder.exists():
        for doc_file in documents_folder.iterdir():
            if doc_file.is_file():
                documents.append({
                    "filename": doc_file.name,
                    "size_bytes": doc_file.stat().st_size,
                    "size_mb": round(doc_file.stat().st_size / (1024 * 1024), 2),
                    "modified": datetime.fromtimestamp(doc_file.stat().st_mtime).isoformat()
                })
    
    return {
        **subjects_db[subject_id],
        "folder_path": str(subject_folder),
        "documents_on_disk": documents
    }

@app.delete("/subjects/{subject_id}/documents/{filename}")
async def delete_document(subject_id: str, filename: str):
    """Delete a specific document from a subject"""
    if subject_id not in subjects_db:
        raise HTTPException(status_code=404, detail="Subject not found")
    
    subject_folder = get_subject_folder(subject_id)
    documents_folder = subject_folder / "documents"
    
    # Find and delete the file
    deleted = False
    for doc_file in documents_folder.iterdir():
        if filename in doc_file.name:
            doc_file.unlink()
            deleted = True
            print(f"✓ Deleted document: {doc_file.name}")
            break
    
    if not deleted:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Note: This doesn't update the vector database
    # For full implementation, you'd need to rebuild the vector DB
    
    return {"message": f"Document deleted: {filename}", "note": "Vector database not updated - consider re-uploading all documents"}

@app.delete("/subjects/{subject_id}")
async def delete_subject(subject_id: str):
    """Delete an entire subject and all its data"""
    if subject_id not in subjects_db:
        raise HTTPException(status_code=404, detail="Subject not found")
    
    import shutil
    subject_folder = get_subject_folder(subject_id)
    
    # Delete entire folder
    shutil.rmtree(subject_folder)
    
    # Remove from memory
    del subjects_db[subject_id]
    
    print(f"✓ Deleted subject and all data: {subject_folder}")
    
    return {"message": "Subject deleted successfully", "deleted_folder": str(subject_folder)}

@app.get("/")
async def root():
    return {
        "message": "RAG Chatbot API is running",
        "model": "HuggingFace Mistral-7B",
        "storage_structure": {
            "base": "rag_data/",
            "subjects": "rag_data/subjects/{subject_id}/",
            "documents": "rag_data/subjects/{subject_id}/documents/",
            "vector_db": "rag_data/subjects/{subject_id}/vector_db/",
            "metadata": "rag_data/subjects/{subject_id}/metadata/"
        }
    }

@app.get("/health")
async def health():
    return {
        "status": "healthy",
        "embedding_model": "all-MiniLM-L6-v2",
        "llm_model": "HuggingFace Mistral-7B Instruct",
        "subjects_count": len(subjects_db),
        "hf_token_set": bool(HF_API_TOKEN)
    }

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)